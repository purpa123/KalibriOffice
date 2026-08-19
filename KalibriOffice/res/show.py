# C:\Users\georg\Projects\CalibriOffice\res\show.py

import customtkinter as ctk
from tkinter import filedialog, messagebox, colorchooser, Menu
import os
import sys

try:
    from PIL import Image, ImageTk
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("⚠️ For images install: pip install Pillow")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ For presentations install: pip install python-pptx")

ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

class KalibriShow(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("KalibriOffice — KalibriShow")
        self.geometry("1100x750")
        self.minsize(800, 500)

        self.current_file = None
        self.slides = []
        self.slide_colors = []
        self.slide_images = []
        self.current_slide_index = 0

        self.create_menu()
        self.create_slide_viewer()
        self.create_controls()
        self.add_slide("Welcome to KalibriShow!")

        print("🎨 KalibriShow started (English UI)")

    def create_menu(self):
        menu_frame = ctk.CTkFrame(self, height=50, fg_color="transparent")
        menu_frame.pack(fill="x", padx=10, pady=(10, 0))

        ctk.CTkButton(menu_frame, text="➕ Add Slide", command=self.add_slide, width=130).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="➖ Delete Slide", command=self.delete_slide, width=130).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="◀️ Previous", command=self.prev_slide, width=80).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="▶️ Next", command=self.next_slide, width=80).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="🎨 Background Color", command=self.choose_background, width=140).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="🖼️ Insert Image", command=self.insert_image, width=160).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="📂 Open", command=self.open_file, width=80).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="💾 Save", command=self.save_file, width=80).pack(side="left", padx=2)
        ctk.CTkButton(menu_frame, text="📁 Save As", command=self.save_as_file, width=110).pack(side="left", padx=2)

        self.label_info = ctk.CTkLabel(menu_frame, text="Slide 1 of 1")
        self.label_info.pack(side="right", padx=10)

    def create_slide_viewer(self):
        self.viewer_frame = ctk.CTkFrame(self)
        self.viewer_frame.pack(padx=10, pady=(10, 5), fill="both", expand=True)

        self.text_area = ctk.CTkTextbox(self.viewer_frame, wrap="word", font=("Arial", 20), height=400)
        self.text_area.pack(fill="both", expand=True, padx=10, pady=10)
        self.text_area.bind("<KeyRelease>", self.on_text_change)
        self.text_area._textbox.bind("<Button-3>", self.show_image_context_menu)

    def create_controls(self):
        control_frame = ctk.CTkFrame(self, height=40, fg_color="transparent")
        control_frame.pack(fill="x", padx=10, pady=(0, 10))

        self.slide_counter = ctk.CTkLabel(control_frame, text="Slide 1 / 1", font=("Arial", 14))
        self.slide_counter.pack(side="left", padx=10)

        ctk.CTkButton(control_frame, text="▶️ Start Slideshow", command=self.start_slideshow, width=140).pack(side="right", padx=10)

    def add_slide(self, content=""):
        self.slides.append(content if content else f"Slide {len(self.slides) + 1}")
        self.slide_colors.append("#2b2b2b")
        self.slide_images.append([])
        self.current_slide_index = len(self.slides) - 1
        self.update_viewer()
        print(f"➕ Added slide {len(self.slides)}")

    def delete_slide(self):
        if len(self.slides) <= 1:
            messagebox.showinfo("KalibriShow", "Cannot delete the last slide.")
            return
        del self.slides[self.current_slide_index]
        del self.slide_colors[self.current_slide_index]
        del self.slide_images[self.current_slide_index]
        if self.current_slide_index >= len(self.slides):
            self.current_slide_index = len(self.slides) - 1
        self.update_viewer()
        print(f"➖ Deleted slide {self.current_slide_index + 1}")

    def next_slide(self):
        if self.current_slide_index < len(self.slides) - 1:
            self.slides[self.current_slide_index] = self.text_area.get("1.0", "end-1c")
            self.current_slide_index += 1
            self.update_viewer()
        else:
            messagebox.showinfo("KalibriShow", "This is the last slide")

    def prev_slide(self):
        if self.current_slide_index > 0:
            self.slides[self.current_slide_index] = self.text_area.get("1.0", "end-1c")
            self.current_slide_index -= 1
            self.update_viewer()
        else:
            messagebox.showinfo("KalibriShow", "This is the first slide")

    def update_viewer(self):
        if self.slides and self.current_slide_index < len(self.slides):
            self.text_area.delete("1.0", "end")
            self.text_area.insert("1.0", self.slides[self.current_slide_index])
            color = self.slide_colors[self.current_slide_index]
            self.text_area.configure(fg_color=color)
            self.label_info.configure(text=f"Slide {self.current_slide_index + 1} of {len(self.slides)}")
            self.slide_counter.configure(text=f"Slide {self.current_slide_index + 1} / {len(self.slides)}")

    def on_text_change(self, event):
        if self.slides and self.current_slide_index < len(self.slides):
            self.slides[self.current_slide_index] = self.text_area.get("1.0", "end-1c")

    def choose_background(self):
        rgb, hex_color = colorchooser.askcolor(
            title="Select background color for this slide",
            initialcolor=self.slide_colors[self.current_slide_index]
        )
        if hex_color:
            self.slide_colors[self.current_slide_index] = hex_color
            self.text_area.configure(fg_color=hex_color)
            print(f"🎨 Slide {self.current_slide_index+1} background changed to {hex_color}")

    def insert_image(self):
        if not PIL_AVAILABLE:
            messagebox.showerror("Error", "Pillow not installed.\nInstall: pip install Pillow")
            return

        file_path = filedialog.askopenfilename(
            title="Select image",
            filetypes=[("Images", "*.png *.jpg *.jpeg *.gif *.bmp")]
        )
        if not file_path:
            return

        try:
            image = Image.open(file_path)
            max_width = 500
            if image.width > max_width:
                ratio = max_width / image.width
                new_size = (max_width, int(image.height * ratio))
                image = image.resize(new_size, Image.Resampling.LANCZOS)
            photo = ImageTk.PhotoImage(image)
            self.text_area._textbox.image_create("insert", image=photo)
            self.text_area._textbox.insert("insert", "\n")
            self.slide_images[self.current_slide_index].append(photo)
            print(f"🖼️ Inserted image into slide {self.current_slide_index + 1}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to insert image:\n{str(e)}")

    def show_image_context_menu(self, event):
        try:
            index = self.text_area._textbox.index(f"@{event.x},{event.y}")
            image_obj = self.text_area._textbox.image_cget(index, "image")
            if not image_obj:
                return
            menu = Menu(self, tearoff=0)
            menu.add_command(
                label="🗑️ Delete image",
                command=lambda: self.delete_image_at_index(index)
            )
            menu.add_separator()
            menu.add_command(label="Cancel", command=lambda: None)
            menu.tk_popup(event.x_root, event.y_root)
        except Exception as e:
            print(f"⚠️ Right-click error: {e}")

    def delete_image_at_index(self, index):
        try:
            self.text_area._textbox.delete(index)
            self.slide_images[self.current_slide_index].clear()
            self.update_viewer()
            print(f"🗑️ Image deleted at {index}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to delete image:\n{str(e)}")

    def start_slideshow(self):
        if not self.slides:
            messagebox.showinfo("KalibriShow", "No slides to show")
            return

        self.slides[self.current_slide_index] = self.text_area.get("1.0", "end-1c")

        show_window = ctk.CTkToplevel(self)
        show_window.title("Slideshow")
        show_window.attributes("-fullscreen", True)
        show_window.grab_set()
        show_window.focus_set()

        slide_frame = ctk.CTkFrame(show_window, fg_color=self.slide_colors[self.current_slide_index])
        slide_frame.pack(fill="both", expand=True)

        label = ctk.CTkLabel(
            slide_frame,
            text=self.slides[self.current_slide_index],
            font=("Arial", 48),
            wraplength=800,
            text_color="white" if self.slide_colors[self.current_slide_index] != "#ffffff" else "black"
        )
        label.pack(expand=True, padx=40, pady=40)

        progress_label = ctk.CTkLabel(
            show_window,
            text=f"{self.current_slide_index + 1} / {len(self.slides)}",
            font=("Arial", 16),
            text_color="gray60"
        )
        progress_label.pack(side="bottom", pady=20)

        def update_fullscreen_slide():
            label.configure(text=self.slides[self.current_slide_index])
            progress_label.configure(text=f"{self.current_slide_index + 1} / {len(self.slides)}")
            color = self.slide_colors[self.current_slide_index]
            slide_frame.configure(fg_color=color)
            label.configure(text_color="white" if color != "#ffffff" else "black")

        def next_slide_fullscreen():
            if self.current_slide_index < len(self.slides) - 1:
                self.current_slide_index += 1
                update_fullscreen_slide()
            else:
                messagebox.showinfo("Slideshow", "End of presentation")

        def prev_slide_fullscreen():
            if self.current_slide_index > 0:
                self.current_slide_index -= 1
                update_fullscreen_slide()

        def close_slideshow():
            show_window.destroy()
            self.update_viewer()

        show_window.bind("<Right>", lambda e: next_slide_fullscreen())
        show_window.bind("<Left>", lambda e: prev_slide_fullscreen())
        show_window.bind("<Escape>", lambda e: close_slideshow())
        show_window.bind("<space>", lambda e: next_slide_fullscreen())

        btn_frame = ctk.CTkFrame(show_window, fg_color="transparent")
        btn_frame.pack(side="bottom", pady=10)

        ctk.CTkButton(btn_frame, text="◀️ Previous", command=prev_slide_fullscreen, width=100).pack(side="left", padx=10)
        ctk.CTkButton(btn_frame, text="▶️ Next", command=next_slide_fullscreen, width=100).pack(side="left", padx=10)
        ctk.CTkButton(btn_frame, text="❌ Exit (Esc)", command=close_slideshow, width=120, fg_color="red").pack(side="left", padx=10)

    def open_file(self):
        if not PPTX_AVAILABLE:
            messagebox.showerror("Error", "python-pptx not installed.\nInstall: pip install python-pptx")
            return

        file_path = filedialog.askopenfilename(
            title="Open presentation",
            filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")]
        )
        if not file_path:
            return

        try:
            prs = Presentation(file_path)
            self.slides = []
            self.slide_colors = []
            self.slide_images = []

            for slide in prs.slides:
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                if not text.strip():
                    text = "Empty slide"
                self.slides.append(text.strip())
                self.slide_colors.append("#2b2b2b")
                self.slide_images.append([])

            if not self.slides:
                self.slides = ["New slide"]
                self.slide_colors = ["#2b2b2b"]
                self.slide_images = [[]]

            self.current_slide_index = 0
            self.current_file = file_path
            self.update_viewer()
            print(f"📂 Opened: {file_path}")

        except Exception as e:
            messagebox.showerror("Error", f"Failed to open:\n{str(e)}")

    def save_file(self):
        if self.current_file:
            self._save_to_file(self.current_file)
        else:
            self.save_as_file()

    def save_as_file(self):
        file_path = filedialog.asksaveasfilename(
            title="Save presentation",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")]
        )
        if not file_path:
            return
        self._save_to_file(file_path)

    def _save_to_file(self, file_path):
        if not PPTX_AVAILABLE:
            messagebox.showerror("Error", "python-pptx not installed.")
            return

        try:
            self.slides[self.current_slide_index] = self.text_area.get("1.0", "end-1c")

            prs = Presentation()
            for idx, content in enumerate(self.slides):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                title = slide.shapes.title
                if title:
                    title.text = f"Slide {idx+1}"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape != title:
                        shape.text = content

            prs.save(file_path)
            self.current_file = file_path
            self.update_title()
            print(f"💾 Saved: {file_path}")

        except Exception as e:
            messagebox.showerror("Error", f"Failed to save:\n{str(e)}")

    def update_title(self):
        if self.current_file:
            name = os.path.basename(self.current_file)
            self.title(f"KalibriShow — {name}")
            self.label_info.configure(text=name)
        else:
            self.title("KalibriShow — Untitled")
            self.label_info.configure(text="Untitled")

if __name__ == "__main__":
    app = KalibriShow()
    app.mainloop()